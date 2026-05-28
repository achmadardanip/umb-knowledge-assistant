from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass
class ExtractedSlide:
    slide_number: int
    slide_title: str | None
    content: str
    extraction_method: str = "python-pptx"
    extraction_confidence: float = 0.8


def extract_pptx(path: str | Path) -> list[ExtractedSlide]:
    try:
        from pptx import Presentation
    except Exception:
        return []

    presentation = Presentation(str(path))
    slides: list[ExtractedSlide] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(
                ExtractedSlide(
                    slide_number=index,
                    slide_title=texts[0][:200] if texts else None,
                    content="\n".join(texts),
                )
            )
    return slides

